#!/usr/bin/env python3
"""
ingest_customer.py  ·  GSM Outdoors — Customer-Facing Wiki Builder
────────────────────────────────────────────────────────────────────
Transforms agent-facing source documents into direct customer-facing
wiki pages, organised brand-wise within each category folder.

WIKI LAYOUT
  wiki/
    index.md                                  ← auto-built master navigation
    fishing/
      fishing-overview.md                     ← category index + cross-brand FAQ
      phenix-rods.md                          ← Phenix brand hub + warranty (single brand)
      phenix-rods-saltwater-pelagic.md        ← Abyss, Abyss HD, Axis, Black Diamond, Pandora
      phenix-rods-saltwater-inshore.md        ← BD Inshore, BD Surf, M1 Inshore, RTS, Titan, Megalodon
      phenix-rods-freshwater-bass.md          ← M1, Feather, Ultra MBX, K2, Ultra Swimbait, BFS
      phenix-rods-freshwater-specialty.md     ← Crankbait, Black Chrome, Maxim, Recon, Kokanee, Walleye
      phenix-rods-trout-ultralight.md         ← Iron Feather, Elixir, Mirage, Dragonfly
      phenix-rods-salmon-steelhead.md         ← Cicada, Trifecta Pro/Lite, Trifecta
      phenix-rods-travel.md                   ← RedEye Saltwater, Freshwater, Trout/Fly
      phenix-rods-warranty.md                 ← Warranty procedures (No-Hassle, Defect, Boron)
      phenix-rods-tier-fees.md                ← Replacement tier fees + return shipping
      phenix-rods-pricing.md                  ← MSRP pricing reference (scraped)
      dobyns-rods.md                          ← Dobyns brand page
      bonehead-tackle.md                      ← Bonehead Tackle brand page
      bucca-brand.md                          ← Bucca Brand page
      dealer-inquiry.md                       ← wholesale / dealer portal
    hunting/
      hunting-overview.md                     ← category index + cross-brand FAQ
      feeders-and-timers.md                   ← WGI + Boss Buck + TH-series (multi-brand)
      avian-x.md                              ← Avian-X brand page
      sog-knives.md                           ← SOG Knives brand page
      replacement-parts.md                    ← Muddy, Hawk, Bloodsport (parts-focused)
      walkers.md                              ← Walker's ear protection brand page
      box-blinds.md                           ← box blinds support
      product-comparisons.md                  ← cross-brand comparison guide
      procedures.md                           ← Bass Pro / Cabela's order procedures
    wireless/
      wireless-overview.md                    ← category index + shared account/app FAQ
      connect-cellular.md                     ← Connect Cellular brand page
      muddy-mtrx.md                           ← Muddy MTRX brand page
      stealth-cam.md                          ← Stealth Cam STC-DS4KTM brand page

RULE: one brand's facts must never appear on another brand's page.
      Overview pages hold cross-brand FAQ only; brand pages hold
      exact procedures, model numbers, addresses, and part lists.

Usage:
    python ingest_customer.py --all
    python ingest_customer.py --all --force
    python ingest_customer.py --file "raw/CS FISHING/Rods - Phenix/Phenix Q&A.docx"
    python ingest_customer.py --folder "raw/CS HUNTING/Feeders & Timers"
    python ingest_customer.py --lint
    python ingest_customer.py --reindex
"""

from __future__ import annotations

import os
import re
import sys
import json
import hashlib
import argparse
from pathlib import Path
from datetime import date

import httpx
from loguru import logger
from rich.console import Console
from rich.table import Table
from rich.progress import track

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE_DIR    = Path(__file__).resolve().parent
RAW_DIR     = BASE_DIR / "raw"
WIKI_DIR    = BASE_DIR / "wiki"
SCHEMA_PATH = BASE_DIR / "SCHEMA.md"
INDEX_PATH  = WIKI_DIR / "index.md"
LOG_PATH    = WIKI_DIR / "log.md"
STATE_PATH  = BASE_DIR / ".ingest_customer_state.json"

# ── Model ─────────────────────────────────────────────────────────────────────
OLLAMA_URL = "http://localhost:11434"
LLM_MODEL  = "gemma4:e4b"

console = Console()

# ── Skip rules ────────────────────────────────────────────────────────────────
SKIP_PATTERNS = [
    "~$", ".DS_Store", "OLD-DONOTUSE", ".onepkg", ".zip", "lint-report",
]
SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt", ".txt", ".md",
}

# ─────────────────────────────────────────────────────────────────────────────
# WIKI PAGE REGISTRY
#
# Each entry defines one wiki page. Key fields:
#   rel_path    : path relative to wiki/
#   title       : page title that appears in frontmatter and index
#   category    : fishing | hunting | wireless | operations
#   brands      : list of brand names this page covers
#   topics      : retrieval keywords (used by wiki_retrieval.py scoring)
#   scope       : one-sentence description fed to the LLM as a guardrail
#   related     : sibling pages to surface alongside this one
#
# Brand-isolation rule: if a page lists exactly one brand, the LLM prompt
# will hard-prohibit content from any other brand. Multi-brand pages (e.g.
# feeders-and-timers.md) explicitly list all brands they may contain.
# ─────────────────────────────────────────────────────────────────────────────

WIKI_PAGES: dict[str, dict] = {


    # ── Fishing — Overview ────────────────────────────────────────────────────
    "fishing/fishing-overview.md": {
        "title":    "Fishing — Product Support Overview",
        "category": "fishing",
        "brands":   ["Phenix Rods", "Dobyns Rods", "Bonehead Tackle", "Bucca Brand"],
        "topics":   ["fishing", "rod", "reel", "lure", "bait", "trout", "bass",
                     "tip", "warranty", "return", "repair"],
        "scope":    (
            "Category overview page for all GSM Outdoors fishing products. "
            "Contains cross-brand FAQ (e.g. 'which fishing brands does GSM carry?') "
            "and short summaries linking out to brand pages. "
            "NEVER contains step-by-step procedures — those belong on brand pages."
        ),
        "related":  ["fishing/phenix-rods.md", "fishing/dobyns-rods.md",
                     "fishing/bonehead-tackle.md", "fishing/bucca-brand.md"],
    },

    # ── Fishing — Phenix Rods (brand hub + warranty) ──────────────────────────
    "fishing/phenix-rods.md": {
        "title":    "Phenix Rods — Brand Overview & Support",
        "category": "fishing",
        "brands":   ["Phenix Rods"],
        "topics":   ["phenix", "phenix rods", "custom quality rods", "gsm outdoors fishing",
                     "phenix warranty", "no-hassle", "phenix repair", "phenix broken rod",
                     "phenix rod tip", "phx", "phenixrods.com"],
        "scope":    (
            "Phenix Rods brand hub. Covers: brand overview, full product category "
            "navigation (saltwater, freshwater, trout, salmon, travel), warranty "
            "program summary, and contact info. Links to sub-pages for detail. "
            "Must NOT contain any content about Dobyns, Bonehead, or any other brand."
        ),
        "related":  [
            "fishing/phenix-rods-warranty.md",
            "fishing/phenix-rods-saltwater-pelagic.md",
            "fishing/phenix-rods-freshwater-bass.md",
            "fishing/phenix-rods-salmon-steelhead.md",
            "fishing/fishing-overview.md",
        ],
    },

    # ── Fishing — Phenix Saltwater ────────────────────────────────────────────
    "fishing/phenix-rods-saltwater-pelagic.md": {
        "title":    "Phenix Rods — Saltwater Pelagic Series",
        "category": "fishing",
        "brands":   ["Phenix Rods"],
        "topics":   ["phenix saltwater", "abyss", "abyss hd", "axis", "black diamond",
                     "black diamond hybrid", "black diamond east coast", "pandora",
                     "live bait rod", "surface jig rod", "pelagic rod", "tuna rod",
                     "yellowtail rod", "west coast fishing", "saltwater casting",
                     "saltwater spinning", "kelp fishing", "phx-psx", "phx-ahd",
                     "phx-hax", "phx-psw", "phx-phd", "phx-esw", "phx-pdx"],
        "scope":    (
            "Phenix Rods saltwater pelagic series ONLY: Abyss, Abyss HD, Axis, "
            "Black Diamond, Black Diamond Hybrid, Black Diamond East Coast, and Pandora. "
            "Covers series descriptions, specs, MSRP, and Best For guidance. "
            "Must NOT contain content about Dobyns, Bonehead, or any other brand."
        ),
        "related":  [
            "fishing/phenix-rods-saltwater-inshore.md",
            "fishing/phenix-rods.md",
            "fishing/phenix-rods-warranty.md",
            "fishing/phenix-rods-tier-fees.md",
        ],
    },
    "fishing/phenix-rods-saltwater-inshore.md": {
        "title":    "Phenix Rods — Saltwater Inshore, Jigging & Surf",
        "category": "fishing",
        "brands":   ["Phenix Rods"],
        "topics":   ["phenix inshore", "m1 inshore", "rts inshore", "rts",
                     "reds trout snook", "black diamond inshore", "black diamond surf",
                     "titan slow jigging", "titan popping", "megalodon",
                     "slow pitch jigging", "long fall jigging", "popping rod",
                     "jigging rod", "surf rod", "surf casting", "gulf coast",
                     "inshore fishing", "phx-smx", "phx-rts", "phx-tjx",
                     "phx-mpx", "phx-tpx", "phx-sx",
                     "tuna", "tuna rod", "tuna jigging", "tuna recommendation"],
        "scope":    (
            "Phenix Rods inshore, jigging, and surf series ONLY: M1 Inshore, "
            "RTS Inshore, Black Diamond Inshore, Black Diamond Surf, Titan Slow Jigging, "
            "Titan Popping, and Megalodon. Covers descriptions, specs, MSRP, Best For. "
            "Must NOT contain content about Dobyns, Bonehead, or any other brand."
        ),
        "related":  [
            "fishing/phenix-rods-saltwater-pelagic.md",
            "fishing/phenix-rods.md",
            "fishing/phenix-rods-warranty.md",
            "fishing/phenix-rods-tier-fees.md",
        ],
    },

    # ── Fishing — Phenix Freshwater ───────────────────────────────────────────
    "fishing/phenix-rods-freshwater-bass.md": {
        "title":    "Phenix Rods — Freshwater Bass (Premium Series)",
        "category": "fishing",
        "brands":   ["Phenix Rods"],
        "topics":   ["phenix bass rod", "m1 bass", "m1 rod", "feather rod", "feather", "m1",
                     "ultra mbx", "ultra mbx classic", "k2 rod", "k2",
                     "ultra swimbait", "ultra swimbait classic", "classic bfs",
                     "bfs rod", "bait finesse", "swimbait rod", "big bait",
                     "tournament bass", "bass rod", "phenix freshwater",
                     "phx-mx", "phx-ftx", "phx-umbxcl", "phx-usb", "phx-tx",
                     "phx-bfs", "nanolite", "toray carbon bass"],
        "scope":    (
            "Phenix Rods premium freshwater bass series ONLY: M1 Bass, Feather, "
            "Ultra MBX Classic, K2, Ultra Swimbait, Ultra Swimbait Classic, and Classic BFS. "
            "Covers descriptions, specs, MSRP, and Best For guidance. "
            "Must NOT contain content about Dobyns, Bonehead, or any other brand."
        ),
        "related":  [
            "fishing/phenix-rods-freshwater-specialty.md",
            "fishing/phenix-rods.md",
            "fishing/phenix-rods-warranty.md",
            "fishing/phenix-rods-tier-fees.md",
        ],
    },
    "fishing/phenix-rods-freshwater-specialty.md": {
        "title":    "Phenix Rods — Freshwater Specialty & Value Series",
        "category": "fishing",
        "brands":   ["Phenix Rods"],
        "topics":   ["maxim", "recon elite", "crankbait xg", "crankbait glass",
                     "crankbait composite x", "composite x", "black chrome",
                     "super flipper", "kokanee reaper", "kokanee rod",
                     "m1 walleye", "walleye rod", "crankbait rod", "glass rod",
                     "drift rod", "value bass rod", "phenix crankbait",
                     "phx-max", "phx-phx", "phx-xg", "phx-x9", "phx-x10",
                     "phx-krx", "phx-wmx"],
        "scope":    (
            "Phenix Rods specialty and value freshwater series ONLY: Maxim, Recon Elite, "
            "Crankbait XG (glass), Crankbait Composite X, Black Chrome, Super Flipper, "
            "Kokanee Reaper, and M1 Walleye. Covers descriptions, specs, MSRP, Best For. "
            "Must NOT contain content about Dobyns, Bonehead, or any other brand."
        ),
        "related":  [
            "fishing/phenix-rods-freshwater-bass.md",
            "fishing/phenix-rods.md",
            "fishing/phenix-rods-warranty.md",
            "fishing/phenix-rods-tier-fees.md",
        ],
    },

    # ── Fishing — Phenix Trout & Ultralight ──────────────────────────────────
    "fishing/phenix-rods-trout-ultralight.md": {
        "title":    "Phenix Rods — Trout & Ultralight Series",
        "category": "fishing",
        "brands":   ["Phenix Rods"],
        "topics":   ["phenix trout", "iron feather", "elixir", "mirage", "dragonfly",
                     "ultralight rod", "trout rod", "finesse rod", "ultralight spinning",
                     "panfish rod", "crappie rod", "1-9 lb", "1-5 lb",
                     "telescopic rod", "burl wood reel seat", "split grip",
                     "phx-if", "phx-fx", "phx-mf", "essex sic trout",
                     "fuji titanium torzile"],
        "scope":    (
            "Phenix Rods trout and ultralight series ONLY: Iron Feather, Elixir, "
            "Mirage, and Dragonfly. Covers descriptions, specs, MSRP, Best For, "
            "and how to choose between the four series. "
            "Must NOT contain content about Dobyns, Bonehead, or any other brand."
        ),
        "related":  [
            "fishing/phenix-rods.md",
            "fishing/phenix-rods-warranty.md",
            "fishing/phenix-rods-tier-fees.md",
        ],
    },

    # ── Fishing — Phenix Salmon & Steelhead ──────────────────────────────────
    "fishing/phenix-rods-salmon-steelhead.md": {
        "title":    "Phenix Rods — Salmon & Steelhead Series",
        "category": "fishing",
        "brands":   ["Phenix Rods"],
        "topics":   ["phenix salmon", "phenix steelhead", "cicada", "trifecta",
                     "trifecta pro", "trifecta lite", "salmon rod", "steelhead rod",
                     "9 ft salmon rod", "9 foot rod", "9 feet", "9 foot salmon", "best salmon rod",
                     "drift fishing", "river fishing",
                     "pacific northwest fishing", "moderate fast salmon",
                     "2 piece salmon", "phx-cax", "phx-tre", "phx-trl", "phx-trx",
                     "best salmon rod phenix", "lightest salmon rod"],
        "scope":    (
            "Phenix Rods salmon and steelhead series ONLY: Cicada, Trifecta Pro, "
            "Trifecta Lite, and Trifecta. Covers descriptions, specs, MSRP, Best For, "
            "and a comparison table for choosing the right 9-foot salmon rod. "
            "Must NOT contain content about Dobyns, Bonehead, or any other brand."
        ),
        "related":  [
            "fishing/phenix-rods.md",
            "fishing/phenix-rods-warranty.md",
            "fishing/phenix-rods-tier-fees.md",
        ],
    },

    # ── Fishing — Phenix Travel ───────────────────────────────────────────────
    "fishing/phenix-rods-travel.md": {
        "title":    "Phenix Rods — Travel Series (RedEye)",
        "category": "fishing",
        "brands":   ["Phenix Rods"],
        "topics":   ["phenix travel rod", "redeye", "red eye travel", "redeye saltwater",
                     "redeye freshwater", "redeye trout", "redeye fly",
                     "multi-piece rod", "backpacking rod", "3-piece rod", "4-piece rod",
                     "travel fishing rod", "hard case rod", "international fishing",
                     "airplane fishing rod", "portable rod", "felt lined case",
                     "phx-rtx", "phx-br", "phx-tr"],
        "scope":    (
            "Phenix Rods RedEye travel series ONLY: RedEye Saltwater, RedEye Freshwater "
            "(Bass/Brass), and RedEye Trout/Fly. Covers descriptions, specs, MSRP, "
            "Best For, and travel-specific questions (airline carry-on, hard case). "
            "Must NOT contain content about Dobyns, Bonehead, or any other brand."
        ),
        "related":  [
            "fishing/phenix-rods.md",
            "fishing/phenix-rods-warranty.md",
            "fishing/phenix-rods-tier-fees.md",
        ],
    },

    # ── Fishing — Phenix Warranty ─────────────────────────────────────────────
    "fishing/phenix-rods-warranty.md": {
        "title":    "Phenix Rods — Warranty & Replacement Programs",
        "category": "fishing",
        "brands":   ["Phenix Rods"],
        "topics":   ["phenix warranty", "no-hassle replacement", "america no-hassle",
                     "broken rod phenix", "rod replacement", "defect claim",
                     "manufacturer defect", "warranty claim", "logo section",
                     "boron legacy", "upgrade rod", "warranty form",
                     "cut rod section", "fishing returns irving tx",
                     "gsm outdoors warranty", "30 days purchase defect",
                     "877-269-8490", "fishinginfo@gsmorg.com",
                     "5250 frye rd", "irving tx 75061",
                     "warranty", "ship rod", "dealer", "snapped", "rod broke",
                     "broken", "process warranty", "warranty through dealer",
                     "upgrade warranty", "rod snapped", "warranty steps"],
        "scope":    (
            "Phenix Rods warranty procedures ONLY. Covers: America No-Hassle "
            "Replacement (steps, logo section cutting, mailing address, payment), "
            "Manufacturer Defect Claim (photo submission, 30-day rule), and Boron "
            "Legacy Program. Exact steps, addresses, form names kept verbatim. "
            "Must NOT contain content about Dobyns, Bonehead, or any other brand."
        ),
        "related":  [
            "fishing/phenix-rods-tier-fees.md",
            "fishing/phenix-rods.md",
        ],
    },

    # ── Fishing — Phenix Tier Fees ────────────────────────────────────────────
    "fishing/phenix-rods-tier-fees.md": {
        "title":    "Phenix Rods — Replacement Tier Fees & Shipping",
        "category": "fishing",
        "brands":   ["Phenix Rods"],
        "topics":   ["phenix tier fee", "replacement fee", "warranty cost",
                     "no-hassle fee", "rod replacement cost", "return shipping",
                     "tier 1", "tier 2", "tier 3", "tier 4", "tier 5",
                     "blank replacement fee", "section replacement fee",
                     "shipping fee warranty", "how much phenix warranty",
                     "upgrade fee calculation", "phenix replacement price",
                     "abyss fee", "black diamond fee", "feather fee",
                     "cicada fee", "trifecta fee", "m1 warranty fee",
                     "feather", "iron feather", "cost to replace",
                     "how much", "upgrade fee", "how much does it cost",
                     "replace a phenix", "tier fee"],
        "scope":    (
            "Phenix Rods replacement tier fee tables ONLY. Covers: complete rod "
            "tier fees (Tier 1–5), blank and section tier fees, return shipping fees "
            "by rod length, upgrade fee formula, and quick-reference series lookup. "
            "Must NOT contain content about Dobyns, Bonehead, or any other brand."
        ),
        "related":  [
            "fishing/phenix-rods-warranty.md",
            "fishing/phenix-rods.md",
        ],
    },

    # ── Fishing — Phenix Pricing ──────────────────────────────────────────────
    "fishing/phenix-rods-pricing.md": {
        "title":    "Phenix Rods — MSRP Pricing Reference",
        "category": "fishing",
        "brands":   ["Phenix Rods"],
        "topics":   ["phenix price", "phenix msrp", "how much phenix rod",
                     "phenix rod cost", "phenix retail price", "phenix rod price",
                     "abyss price", "cicada price", "feather price", "m1 price",
                     "black diamond price", "iron feather price", "k2 price",
                     "trifecta price", "redeye price", "axis price",
                     "ultra mbx price", "maxim price", "recon elite price",
                     "cicada", "affordable", "premium", "most affordable",
                     "most premium", "how much", "cicada cost", "cheapest"],
        "scope":    (
            "Phenix Rods MSRP pricing reference ONLY. Contains price ranges for "
            "all series scraped from phenixrods.com. Used to answer 'how much does "
            "X cost?' and to calculate warranty upgrade fees. "
            "Must NOT contain content about Dobyns, Bonehead, or any other brand."
        ),
        "related":  [
            "fishing/phenix-rods-tier-fees.md",
            "fishing/phenix-rods-warranty.md",
            "fishing/phenix-rods.md",
        ],
    },

    # ── Fishing — Other brands ────────────────────────────────────────────────
    "fishing/dobyns-rods.md": {
        "title":    "Dobyns Rods — Warranty & Support",
        "category": "fishing",
        "brands":   ["Dobyns Rods"],
        "topics":   ["dobyns", "champion", "sierra", "dobyns warranty",
                     "dobyns rod", "dobyns repair", "dobyns care"],
        "scope":    (
            "Dobyns Rods ONLY. Covers: warranty process (exact steps, addresses, "
            "form names, fees), rod care instructions, Champion and Sierra series "
            "info, FAQ, model numbers. "
            "Must NOT contain any content about Phenix, Bonehead, or any other brand."
        ),
        "related":  ["fishing/fishing-overview.md", "fishing/phenix-rods.md"],
    },
    "fishing/bonehead-tackle.md": {
        "title":    "Bonehead Tackle — E-Series Rod Support",
        "category": "fishing",
        "brands":   ["Bonehead Tackle"],
        "topics":   ["bonehead", "e-series", "carbon fiber rod", "bonehead tackle",
                     "bonehead warranty", "bonehead bait"],
        "scope":    (
            "Bonehead Tackle ONLY. Covers: E-Series Carbon Fiber rods, baits, "
            "pro staff info, warranty process, FAQ, and part numbers. "
            "Must NOT contain content about any other fishing brand."
        ),
        "related":  ["fishing/fishing-overview.md"],
    },
    "fishing/bucca-brand.md": {
        "title":    "Bucca Brand — Brand Overview",
        "category": "fishing",
        "brands":   ["Bucca Brand"],
        "topics":   ["bucca", "bucca brand", "mike bucca", "trick shad", "baby bull shad",
                     "baby bull gill", "bull mullet", "baby bull rat", "bull wake",
                     "buzzing baby bull shad", "weedless baby bull shad", "swimbait",
                     "glide bait", "ABS swimbait", "bucca swimbaits"],
        "scope":    (
            "Bucca Brand ONLY. High-level brand overview, product families summary, "
            "freshwater and saltwater product lines. Links to catalog, warranty, returns pages."
        ),
        "related":  ["fishing/fishing-overview.md", "fishing/bucca-brand-product-catalog.md",
                     "fishing/bucca-brand-warranty.md", "fishing/bucca-brand-returns.md"],
    },
    "fishing/bucca-brand-product-catalog.md": {
        "title":    "Bucca Brand — Product Catalog",
        "category": "fishing",
        "brands":   ["Bucca Brand"],
        "topics":   ["trick shad", "4 inch trick shad", "6 inch trick shad", "8 inch trick shad",
                     "bfs trick shad", "baby bull shad", "lil baby bull shad", "5 inch bull shad",
                     "baby bull gill", "buzzing baby bull shad", "weedless baby bull shad",
                     "baby bull rat", "bull mullet", "5 inch bull wake", "replacement tail",
                     "bucca brand colors", "bucca brand SKU", "BUC-TS6", "BUC-BBS375",
                     "BUC-BBG375", "BUC-BBR35", "BUC-BM55", "BUC-BM8", "BUC-BBBS375",
                     "BUC-WBBS45", "BUC-BSWB5", "BUC-BS5", "swimbait", "glide bait",
                     "mike bucca", "ABS swimbait"],
        "scope":    (
            "Bucca Brand ONLY. Complete product catalog with model details, SKUs, color tables, "
            "and replacement parts. All freshwater and saltwater products covered."
        ),
        "related":  ["fishing/bucca-brand.md", "fishing/bucca-brand-warranty.md",
                     "fishing/bucca-brand-returns.md"],
    },
    "fishing/bucca-brand-warranty.md": {
        "title":    "Bucca Brand — Warranty Replacement Policy",
        "category": "fishing",
        "brands":   ["Bucca Brand"],
        "topics":   ["bucca brand warranty", "bucca warranty claim", "bucca brand defect",
                     "bucca brand replacement", "bucca warranty", "manufacture defect",
                     "warranty replacement request", "RA number", "return authorization",
                     "bucca brand authorized retailer", "proof of purchase bucca",
                     "21 days replacement"],
        "scope":    (
            "Bucca Brand ONLY. GSM Outdoors warranty replacement policy: eligibility, "
            "claim process, RA number requirement, 21-day timeline, shipping rules."
        ),
        "related":  ["fishing/bucca-brand.md", "fishing/bucca-brand-product-catalog.md",
                     "fishing/bucca-brand-returns.md"],
    },
    "fishing/bucca-brand-returns.md": {
        "title":    "Bucca Brand — Returns & Refunds Policy",
        "category": "fishing",
        "brands":   ["Bucca Brand"],
        "topics":   ["bucca brand return", "bucca brand refund", "bucca return request",
                     "GSM return policy", "return authorization", "RA number",
                     "30 day return", "free return", "bucca brand exchange",
                     "return shipping label", "original packaging", "GSM Outdoors refund",
                     "baits.com return", "21 business days"],
        "scope":    (
            "Bucca Brand ONLY. GSM Outdoors return and refund policy: 30-day window, "
            "eligibility, RA number, pre-paid label, 21 business day refund timeline."
        ),
        "related":  ["fishing/bucca-brand.md", "fishing/bucca-brand-product-catalog.md",
                     "fishing/bucca-brand-warranty.md"],
    },
    "fishing/dealer-inquiry.md": {
        "title":    "Becoming a GSM Outdoors Dealer",
        "category": "operations",
        "brands":   [],
        "topics":   ["dealer", "wholesale", "distributor", "territory",
                     "credit app", "reseller", "retail partner"],
        "scope":    (
            "Wholesale dealer and reseller inquiry process. Covers: how to apply, "
            "territory questions, credit application, and next steps. "
            "Not about consumer warranty or product support."
        ),
        "related":  [],
    },

    # ── Hunting ───────────────────────────────────────────────────────────────
    "hunting/hunting-overview.md": {
        "title":    "Hunting — Product Support Overview",
        "category": "hunting",
        "brands":   ["Avian-X", "Boss Buck", "WGI", "SOG", "Muddy", "Hawk",
                     "Walker's", "Bloodsport"],
        "topics":   ["hunting", "feeder", "decoy", "knife", "blind", "stand",
                     "camera", "ear protection", "replacement part"],
        "scope":    (
            "Category overview for all GSM Outdoors hunting products. "
            "Contains cross-brand FAQ and short summaries linking to brand pages. "
            "NEVER contains step-by-step procedures — those belong on brand pages."
        ),
        "related":  ["hunting/feeders-and-timers.md", "hunting/avian-x.md",
                     "hunting/sog-knives.md", "hunting/replacement-parts.md",
                     "hunting/walkers.md"],
    },
    "hunting/feeders-and-timers.md": {
        "title":    "Wildlife Feeders & Timers — WGI, Boss Buck, TH-Series",
        "category": "hunting",
        "brands":   ["WGI", "Wildgame Innovations", "Boss Buck", "TH-270", "TH-DX1"],
        "topics":   ["feeder", "timer", "boss buck", "wgi", "wildgame", "th-270",
                     "th-dx1", "battery", "spinner", "dispense", "motor", "jam",
                     "short", "error", "feeder error", "feeder battery",
                     "feeder motor", "feeder timer", "ah-series", "xde", "xa"],
        "scope":    (
            "Wildlife feeders and timers sold by GSM Outdoors: WGI (Wildgame Innovations), "
            "Boss Buck, and TH-series (TH-270, TH-DX1). Covers: troubleshooting error codes, "
            "battery and motor issues, spinning mechanism jams, replacement parts, "
            "and warranty procedures — with exact model numbers and steps. "
            "Keep WGI content clearly separated from Boss Buck content using sub-headings."
        ),
        "related":  ["hunting/replacement-parts.md", "hunting/hunting-overview.md"],
    },
    "hunting/avian-x.md": {
        "title":    "Avian-X — Decoys & Blinds Support",
        "category": "hunting",
        "brands":   ["Avian-X"],
        "topics":   ["avian", "avian-x", "decoy", "blind", "lcd valve",
                     "avian-x parts", "avian x", "a-frame blind"],
        "scope":    (
            "Avian-X ONLY. Covers: LCD Mallard and other decoy models, A-frame blinds, "
            "replacement parts list, warranty, FAQ. "
            "Must NOT contain content about any other hunting brand."
        ),
        "related":  ["hunting/replacement-parts.md", "hunting/hunting-overview.md"],
    },
    "hunting/sog-knives.md": {
        "title":    "SOG Knives — Cutlery & Engraving Support",
        "category": "hunting",
        "brands":   ["SOG"],
        "topics":   ["sog", "knife", "knives", "cutlery", "blade",
                     "engraving", "sog warranty", "sog sku"],
        "scope":    (
            "SOG Knives ONLY. Covers: product line overview, SKU information, "
            "engraving service details, warranty, and FAQ. "
            "Must NOT contain content about any other brand."
        ),
        "related":  ["hunting/hunting-overview.md"],
    },
    "hunting/replacement-parts.md": {
        "title":    "Replacement Parts — Muddy, Hawk, Bloodsport & Others",
        "category": "hunting",
        "brands":   ["Muddy", "Hawk", "Bloodsport", "WGI"],
        "topics":   ["muddy", "hawk", "replacement part", "stand", "ladder",
                     "cable", "bloodsport", "outsert", "tree stand part",
                     "ladder stand", "treestand"],
        "scope":    (
            "Replacement parts for Muddy treestands, Hawk ladder stands, "
            "Bloodsport arrows/outserts, and WGI feeder parts. "
            "Covers part numbers, ordering steps, compatibility, and FAQ. "
            "Keep each brand in its own sub-section."
        ),
        "related":  ["hunting/feeders-and-timers.md", "hunting/hunting-overview.md"],
    },
    "hunting/walkers.md": {
        "title":    "Walker's — Ear Protection & Game Ear Support",
        "category": "hunting",
        "brands":   ["Walker's"],
        "topics":   ["walker", "ear protection", "game ear", "hearing",
                     "nrr", "walker's", "walkers game ear", "muffs"],
        "scope":    (
            "Walker's ear protection products ONLY. Covers: Game Ear models, "
            "NRR ratings, battery replacement, troubleshooting, and warranty. "
            "Must NOT contain content about any other brand."
        ),
        "related":  ["hunting/product-comparisons.md", "hunting/hunting-overview.md"],
    },
    "hunting/box-blinds.md": {
        "title":    "Box Blinds — Parts & Assembly",
        "category": "hunting",
        "brands":   [],
        "topics":   ["box blind", "hunting blind", "blind panel", "blind assembly",
                     "blind replacement", "blind part"],
        "scope":    (
            "Box blinds sold by GSM Outdoors. Covers: panel part numbers, "
            "assembly instructions, replacement parts, FAQ."
        ),
        "related":  ["hunting/replacement-parts.md", "hunting/hunting-overview.md"],
    },
    "hunting/product-comparisons.md": {
        "title":    "Product Comparison Guide",
        "category": "hunting",
        "brands":   [],
        "topics":   ["compare", "comparison", "versus", "vs", "difference between",
                     "better", "which is better", "should i get"],
        "scope":    (
            "Side-by-side comparison tables for GSM Outdoors products: game cameras, "
            "box blinds, Walker's ear protection models. Helps customers choose. "
            "References exact model names — do not embed full procedures here."
        ),
        "related":  ["hunting/walkers.md", "wireless/wireless-overview.md"],
    },
    "hunting/procedures.md": {
        "title":    "Bass Pro / Cabela's Order Procedures",
        "category": "operations",
        "brands":   [],
        "topics":   ["bass pro", "cabela", "bc ", "unit of measure",
                     "procedure", "bass pro order", "cabela order"],
        "scope":    (
            "Purchase and order procedures specific to Bass Pro Shops and Cabela's "
            "(BC) customers. Covers: unit-of-measure rules, ordering workflow, "
            "and any BC-specific policies."
        ),
        "related":  [],
    },

    # ── Wireless ──────────────────────────────────────────────────────────────
    "wireless/wireless-overview.md": {
        "title":    "Wireless Cameras — Account Setup & App FAQ",
        "category": "wireless",
        "brands":   ["Connect Cellular", "Muddy MTRX", "Stealth Cam"],
        "topics":   ["camera", "cellular camera", "trail cam", "app", "account",
                     "data plan", "at&t", "verizon", "signal", "solar panel",
                     "connector", "wireless setup", "gsm app"],
        "scope":    (
            "Category overview for all GSM Outdoors wireless cameras. "
            "Contains: shared account setup steps, app FAQ, data plan comparison, "
            "AT&T / Verizon compatibility, and cross-brand FAQ. "
            "Links to brand pages for model-specific troubleshooting. "
            "NEVER contains model-specific steps — those belong on brand pages."
        ),
        "related":  ["wireless/connect-cellular.md", "wireless/muddy-mtrx.md",
                     "wireless/stealth-cam.md"],
    },
    "wireless/connect-cellular.md": {
        "title":    "Connect Cellular Camera — Setup & Troubleshooting",
        "category": "wireless",
        "brands":   ["Connect Cellular"],
        "topics":   ["connect", "connect cellular", "maneuver", "led", "red light",
                     "sd card", "sim card", "antenna", "sync button", "data plan",
                     "connect camera", "maneuver camera"],
        "scope":    (
            "Connect Cellular / Maneuver camera ONLY. Covers: initial setup, "
            "LED status codes, SD card and SIM card issues, sync button procedure, "
            "antenna replacement, data plan activation, troubleshooting steps, "
            "and warranty. Must NOT contain content about MTRX or Stealth Cam."
        ),
        "related":  ["wireless/wireless-overview.md", "wireless/muddy-mtrx.md"],
    },
    "wireless/muddy-mtrx.md": {
        "title":    "Muddy MTRX Camera — Setup & Troubleshooting",
        "category": "wireless",
        "brands":   ["Muddy MTRX"],
        "topics":   ["mtrx", "muddy mtrx", "mud-mtrx", "mtrx camera",
                     "mtrx setup", "mtrx troubleshoot"],
        "scope":    (
            "Muddy MTRX camera ONLY. Covers: setup, app pairing, LED codes, "
            "troubleshooting, and warranty. "
            "Must NOT contain content about Connect Cellular or Stealth Cam."
        ),
        "related":  ["wireless/wireless-overview.md", "wireless/connect-cellular.md"],
    },
    "wireless/stealth-cam.md": {
        "title":    "Stealth Cam STC-DS4KTM — Setup & Troubleshooting",
        "category": "wireless",
        "brands":   ["Stealth Cam"],
        "topics":   ["stealth cam", "stc-ds4k", "ds4ktm", "stealthcam",
                     "stealth cam setup", "stealth cam troubleshoot"],
        "scope":    (
            "Stealth Cam STC-DS4KTM ONLY. Covers: setup, app pairing, "
            "troubleshooting, and warranty. "
            "Must NOT contain content about Connect Cellular or Muddy MTRX."
        ),
        "related":  ["wireless/wireless-overview.md"],
    },
}

# ── Folder → wiki page(s) mapping ─────────────────────────────────────────────
# Order matters: more-specific entries must appear before broader ones.
# The first matching fragment wins unless multiple pages are listed.
#
# Phenix split strategy:
#   - Warranty/Q&A docs → phenix-rods-warranty.md (procedures) + phenix-rods.md (hub)
#   - Catalog/product info → routed to correct product sub-page by content
#   - Tier fee / replacement form → phenix-rods-tier-fees.md
#   - General Phenix docs → phenix-rods.md (hub)

FOLDER_WIKI_MAP: list[tuple[str, list[str]]] = [
    # ── Fishing — Phenix (specific sub-folders first) ─────────────────────────
    ("Phenix Warranty",                    ["fishing/phenix-rods-warranty.md",
                                            "fishing/phenix-rods-tier-fees.md"]),
    ("FISHING ROD WARRANTY FORMS",         ["fishing/phenix-rods-warranty.md",
                                            "fishing/phenix-rods-tier-fees.md",
                                            "fishing/dobyns-rods.md"]),
    ("Phenix Saltwater",                   ["fishing/phenix-rods-saltwater-pelagic.md",
                                            "fishing/phenix-rods-saltwater-inshore.md"]),
    ("Phenix Freshwater",                  ["fishing/phenix-rods-freshwater-bass.md",
                                            "fishing/phenix-rods-freshwater-specialty.md"]),
    ("Phenix Trout",                       ["fishing/phenix-rods-trout-ultralight.md"]),
    ("Phenix Salmon",                      ["fishing/phenix-rods-salmon-steelhead.md"]),
    ("Phenix Travel",                      ["fishing/phenix-rods-travel.md"]),
    # General Phenix folder → hub + warranty (broadest Phenix catch-all, keep last)
    ("Rods - Phenix",                      ["fishing/phenix-rods.md",
                                            "fishing/phenix-rods-warranty.md"]),

    # ── Fishing — Other brands ────────────────────────────────────────────────
    ("Rods and Baits - Bonehead Tackle",   ["fishing/bonehead-tackle.md"]),
    ("Rods and Baits - Dobyns",            ["fishing/dobyns-rods.md"]),
    ("Baits - Bucca Brand",                ["fishing/bucca-brand.md",
                                            "fishing/bucca-brand-product-catalog.md",
                                            "fishing/bucca-brand-warranty.md",
                                            "fishing/bucca-brand-returns.md"]),
    ("Fishing Dealer Inquiry",             ["fishing/dealer-inquiry.md"]),
    ("Fishing Training",                   ["fishing/fishing-overview.md"]),
    ("CS FISHING",                         ["fishing/fishing-overview.md"]),

    # ── Hunting ───────────────────────────────────────────────────────────────
    ("Avian-X Parts",                      ["hunting/avian-x.md",
                                            "hunting/replacement-parts.md"]),
    ("Feeder & Timer Manuals",             ["hunting/feeders-and-timers.md"]),
    ("Feeders & Timers",                   ["hunting/feeders-and-timers.md"]),
    ("Boss Buck",                          ["hunting/feeders-and-timers.md"]),
    ("SOG",                                ["hunting/sog-knives.md"]),
    ("Comparison Docs",                    ["hunting/product-comparisons.md"]),
    ("Procedures",                         ["hunting/procedures.md"]),
    ("Replacement Parts File",             ["hunting/replacement-parts.md"]),
    ("CS HUNTING",                         ["hunting/hunting-overview.md",
                                            "hunting/replacement-parts.md"]),
    ("Walker",                             ["hunting/walkers.md"]),
    ("Box Blind",                          ["hunting/box-blinds.md"]),

    # ── Wireless ──────────────────────────────────────────────────────────────
    ("Connect Cellular",                   ["wireless/connect-cellular.md",
                                            "wireless/wireless-overview.md"]),
    ("MUD-MTRX",                           ["wireless/muddy-mtrx.md",
                                            "wireless/wireless-overview.md"]),
    ("STC-DS4KTM",                         ["wireless/stealth-cam.md",
                                            "wireless/wireless-overview.md"]),
    ("Camera Manuals",                     ["wireless/wireless-overview.md"]),
    ("GT - Wireless Tech Support",         ["wireless/wireless-overview.md"]),

]


# ─────────────────────────────────────────────────────────────────────────────
# Prompt templates
# ─────────────────────────────────────────────────────────────────────────────

REFRAME_INSTRUCTION = """\
=== LANGUAGE TRANSFORMATION RULES ===

The source document is written FOR customer support AGENTS.
Rewrite it FOR customers reading it directly.

  "Tell the customer to..."           → write the step directly
  "Advise the customer to..."         → write the step directly
  "If customer asks, respond with..." → write the answer directly
  "Direct the customer to..."         → write the step directly
  "Inform the customer that..."       → state the fact directly
  "The customer should be told..."    → state the fact directly

ACCURACY — THIS IS THE MOST IMPORTANT RULE:
  Copy specific procedures EXACTLY. Never summarise or replace steps
  with vague phrases like "contact customer service."

  Source: "Have customer cut a 6-inch section containing the rod logo,
           mail it with the Warranty Information Form and $10 to:
           Fishing Returns, 123 Main St, Irving TX 75038"
  Wiki:   "1. Cut a 6-inch section of the rod that contains the logo.
           2. Complete the Warranty Information Form.
           3. Include a $10 payment.
           4. Mail to: Fishing Returns, 123 Main St, Irving TX 75038"

OMIT:
  - Agent escalation notes ("escalate to tier 2", "transfer to supervisor")
  - Internal pricing/campaign codes (e.g. "GTSALES")
  - Scripted agent greetings and closings
  - Internal notes irrelevant to the customer

KEEP VERBATIM:
  - Every step in every procedure
  - Mailing addresses, form names, part numbers
  - URLs and links
  - Model numbers and specs
  - Prices the customer needs to know
  - Warranty terms and timeframes
=== END LANGUAGE TRANSFORMATION RULES ===
"""

# Frontmatter template — filled in by build_frontmatter()
FRONTMATTER_TEMPLATE = """\
---
title: "{title}"
brands: {brands_yaml}
topics: {topics_yaml}
category: {category}
source_docs:
{source_docs_yaml}
last_updated: {today}
related_pages: {related_yaml}
---
"""

CREATE_PROMPT = """\
You are creating a customer-facing support wiki page for GSM Outdoors.
{reframe}

TARGET FILE : wiki/{rel_path}
PAGE TITLE  : {title}
THIS PAGE IS ABOUT: {scope}

BRAND ISOLATION RULE:
{brand_rule}

REQUIRED OUTPUT STRUCTURE (write directly TO customers):
1. Start immediately with the YAML frontmatter block below — copy it exactly,
   do NOT alter any field values:

{frontmatter}

2. Then write the page body using this structure:
   ## Overview
   (One paragraph: what this brand/product is and what customers can get help with here)

   ## [Main topic 1 — e.g. Warranty, Setup, Troubleshooting, Parts]
   (Use numbered lists for multi-step procedures)
   (Use tables for parts: | Part Name | Part Number | Notes |)

   ## [More topic sections as needed]

   ## Common Questions
   **[Question a customer would actually ask?]**
   [Direct answer — exact steps/info from the source]

   ## Sources
   - {source_file}

AGENT-FACING SOURCE TO TRANSFORM:
{raw_text}

Return ONLY the complete wiki page starting with ---"""

UPDATE_PROMPT = """\
You are updating an existing customer-facing support wiki page for GSM Outdoors.
{reframe}

TARGET FILE : wiki/{rel_path}
PAGE TITLE  : {title}
THIS PAGE IS ABOUT: {scope}

BRAND ISOLATION RULE:
{brand_rule}

UPDATE RULES:
1. Merge new information into the correct existing sections — do NOT duplicate.
2. Update the frontmatter: add "{source_file}" to source_docs (if not already there),
   set last_updated to {today}.
3. Preserve all existing content that is still valid.
4. Keep the exact frontmatter field order shown in the existing page.
5. Return the COMPLETE updated page — no preamble or commentary.

EXISTING PAGE:
{existing}

NEW AGENT-FACING SOURCE TO TRANSFORM:
{raw_text}

Return ONLY the complete updated wiki page starting with ---"""


def _brand_rule(page_def: dict) -> str:
    brands = page_def.get("brands", [])
    if not brands:
        return "This page has no brand restriction — include all relevant content."
    if len(brands) == 1:
        return (
            f"This page covers {brands[0]} ONLY.\n"
            f"HARD RULE: Do NOT include any information about other brands.\n"
            f"If the source document mentions other brands, ignore that content."
        )
    brand_list = ", ".join(brands)
    return (
        f"This page covers these brands: {brand_list}.\n"
        f"HARD RULE: Do NOT include information about brands not in that list.\n"
        f"Use clear sub-headings (### {brands[0]}, ### {brands[1]}, …) to separate each brand."
    )


def build_frontmatter(page_def: dict, source_file: str,
                      existing_sources: list[str] | None = None) -> str:
    """Build a complete YAML frontmatter block for a wiki page."""
    today = date.today().isoformat()
    sources = list(existing_sources or [])
    if source_file and source_file not in sources:
        sources.append(source_file)

    def yaml_list(items: list[str]) -> str:
        if not items:
            return "[]"
        quoted = [f'"{x}"' if " " in x or "," in x else x for x in items]
        return "[" + ", ".join(quoted) + "]"

    source_lines = "\n".join(f"  - {s}" for s in sources) if sources else "  []"
    return FRONTMATTER_TEMPLATE.format(
        title=page_def["title"],
        brands_yaml=yaml_list(page_def["brands"]),
        topics_yaml=yaml_list(page_def["topics"]),
        category=page_def["category"],
        source_docs_yaml=source_lines,
        today=today,
        related_yaml=yaml_list(page_def.get("related", [])),
    )


# ─────────────────────────────────────────────────────────────────────────────
# Utility helpers
# ─────────────────────────────────────────────────────────────────────────────

def should_skip(path: Path) -> bool:
    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return True
    name = path.name
    for pat in SKIP_PATTERNS:
        if pat.lower() in name.lower():
            return True
    return False


def file_hash(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def load_state() -> dict:
    if STATE_PATH.exists():
        try:
            return json.loads(STATE_PATH.read_text())
        except Exception:
            return {}
    return {}


def save_state(state: dict) -> None:
    STATE_PATH.write_text(json.dumps(state, indent=2))


def is_unchanged(path: Path, state: dict) -> bool:
    return state.get(str(path)) == file_hash(path)


def mark_done(path: Path, state: dict) -> None:
    state[str(path)] = file_hash(path)


def has_valid_frontmatter(content: str) -> bool:
    return content.strip().startswith("---")


def parse_existing_sources(content: str) -> list[str]:
    """Extract source_docs list from existing frontmatter."""
    if not has_valid_frontmatter(content):
        return []
    end = content.find("\n---\n", 3)
    if end == -1:
        return []
    fm = content[3:end]
    sources: list[str] = []
    in_sources = False
    for line in fm.splitlines():
        if re.match(r"^source_docs:\s*$", line):
            in_sources = True
            continue
        if in_sources:
            m = re.match(r"^\s+-\s+(.+)$", line)
            if m:
                sources.append(m.group(1).strip())
            elif line.strip() and not line.startswith(" "):
                break
    return sources


def resolve_wiki_pages(file_path: Path) -> list[str]:
    """Map a raw file path to one or more target wiki page keys."""
    path_str = str(file_path).replace("\\", "/")
    for fragment, pages in FOLDER_WIKI_MAP:
        if fragment.lower() in path_str.lower():
            return pages
    logger.warning("No wiki mapping for: {}", file_path.name)
    return []


def append_log(entry: str) -> None:
    today = date.today().isoformat()
    LOG_PATH.parent.mkdir(parents=True, exist_ok=True)
    if not LOG_PATH.exists():
        LOG_PATH.write_text("# GSM Outdoors CS Wiki — Operation Log\n\n")
    with open(LOG_PATH, "a") as f:
        f.write(f"\n## [{today}] {entry}\n")


# ─────────────────────────────────────────────────────────────────────────────
# Text extraction
# ─────────────────────────────────────────────────────────────────────────────

def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            return _extract_pdf(path)
        elif ext in (".docx", ".doc"):
            return _extract_docx(path)
        elif ext in (".xlsx", ".xls"):
            return _extract_xlsx(path)
        elif ext in (".pptx", ".ppt"):
            return _extract_pptx(path)
        elif ext in (".txt", ".md"):
            return path.read_text(errors="replace")
    except Exception as e:
        logger.error("Extraction failed for {}: {}", path.name, e)
    return ""


def _extract_pdf(path: Path) -> str:
    try:
        from docling.document_converter import DocumentConverter
        result = DocumentConverter().convert(str(path))
        text = result.document.export_to_markdown()
        if len(text.strip()) > 50:
            return text
    except Exception:
        pass
    try:
        from pdfminer.high_level import extract_text as pdf_extract
        return pdf_extract(str(path)) or ""
    except Exception as e:
        logger.warning("PDF fallback failed for {}: {}", path.name, e)
    return ""


def _extract_docx(path: Path) -> str:
    import docx
    doc = docx.Document(str(path))
    parts = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        rows = []
        for i, row in enumerate(table.rows):
            cells = [c.text.strip() for c in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows.append("| " + " | ".join(["---"] * len(row.cells)) + " |")
        if rows:
            parts.append("\n".join(rows))
    return "\n\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    import pandas as pd
    xl = pd.ExcelFile(str(path))
    parts = []
    for sheet in xl.sheet_names:
        try:
            df = xl.parse(sheet, dtype=str).fillna("")
            if df.empty:
                continue
            parts.append(f"### Sheet: {sheet}\n\n{df.to_markdown(index=False)}")
        except Exception as e:
            logger.warning("Sheet '{}' in {} failed: {}", sheet, path.name, e)
    return "\n\n".join(parts)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        if texts:
            parts.append(f"### Slide {i}\n\n" + "\n".join(texts))
    return "\n\n".join(parts)


# ─────────────────────────────────────────────────────────────────────────────
# LLM call — retry + strip thinking tags
# ─────────────────────────────────────────────────────────────────────────────

def llm(prompt: str, system: str = "", max_tokens: int = 3000) -> str:
    messages = []
    if system:
        messages.append({"role": "system", "content": system})
    messages.append({"role": "user", "content": prompt})

    for attempt in range(1, 4):
        try:
            with httpx.Client(timeout=300) as client:
                resp = client.post(
                    f"{OLLAMA_URL}/api/chat",
                    json={
                        "model":   LLM_MODEL,
                        "messages": messages,
                        "stream":  False,
                        "options": {
                            "temperature": 0.1,
                            "num_predict": max_tokens,
                            "num_ctx":     8192,
                        },
                    },
                )
                resp.raise_for_status()
                raw = resp.json()["message"]["content"].strip()
                # Strip Gemma thinking blocks
                clean = re.sub(r"<think>.*?</think>", "", raw, flags=re.DOTALL).strip()
                # Strip accidental markdown fences wrapping the whole page
                clean = re.sub(r"^```(?:markdown|md)?\s*\n", "", clean)
                clean = re.sub(r"\n```\s*$", "", clean)
                if clean:
                    return clean
                logger.warning("Attempt {}/3 — empty response", attempt)
        except httpx.ConnectError:
            logger.error("Cannot reach Ollama — run: ollama serve")
            sys.exit(1)
        except Exception as e:
            logger.error("LLM error (attempt {}/3): {}", attempt, e)

    logger.error("All 3 LLM attempts failed — skipping")
    return ""


def get_schema() -> str:
    return SCHEMA_PATH.read_text() if SCHEMA_PATH.exists() else ""


# ─────────────────────────────────────────────────────────────────────────────
# Core ingest
# ─────────────────────────────────────────────────────────────────────────────

def ingest_file(file_path: Path, force: bool = False) -> dict:
    stats: dict = {
        "file":          file_path.name,
        "status":        "skipped",
        "pages_updated": [],
        "error":         None,
    }

    if should_skip(file_path):
        return stats

    state = load_state()
    if not force and is_unchanged(file_path, state):
        logger.debug("Unchanged: {}", file_path.name)
        return stats

    logger.info("Ingesting: {}", file_path.name)

    raw_text = extract_text(file_path)
    if not raw_text or len(raw_text.strip()) < 20:
        logger.warning("No usable text: {}", file_path.name)
        stats["status"] = "empty"
        mark_done(file_path, state)
        save_state(state)
        return stats

    # Truncate very large documents to keep within LLM context
    if len(raw_text) > 5000:
        raw_text = raw_text[:5000] + "\n\n[... document truncated for ingest ...]"

    target_pages = resolve_wiki_pages(file_path)
    if not target_pages:
        stats["status"] = "no_mapping"
        return stats

    schema = get_schema()

    for rel_path in target_pages:
        page_def = WIKI_PAGES.get(rel_path)
        if page_def is None:
            logger.warning("rel_path '{}' not in WIKI_PAGES registry — skipping", rel_path)
            continue

        wiki_page_path = WIKI_DIR / rel_path
        wiki_page_path.parent.mkdir(parents=True, exist_ok=True)

        existing_content = ""
        if wiki_page_path.exists():
            existing_content = wiki_page_path.read_text(errors="replace").strip()

        brand_rule = _brand_rule(page_def)

        if existing_content and has_valid_frontmatter(existing_content):
            # ── UPDATE existing page ──────────────────────────────────────
            existing_sources = parse_existing_sources(existing_content)
            prompt = UPDATE_PROMPT.format(
                reframe=REFRAME_INSTRUCTION,
                rel_path=rel_path,
                title=page_def["title"],
                scope=page_def["scope"],
                brand_rule=brand_rule,
                source_file=file_path.name,
                today=date.today().isoformat(),
                existing=existing_content[:2500],
                raw_text=raw_text,
            )
        else:
            # ── CREATE new page ───────────────────────────────────────────
            frontmatter = build_frontmatter(page_def, file_path.name)
            prompt = CREATE_PROMPT.format(
                reframe=REFRAME_INSTRUCTION,
                rel_path=rel_path,
                title=page_def["title"],
                scope=page_def["scope"],
                brand_rule=brand_rule,
                frontmatter=frontmatter,
                source_file=file_path.name,
                raw_text=raw_text,
            )

        result = llm(prompt, system=schema)

        if not result:
            logger.warning("  Empty LLM result for wiki/{}", rel_path)
            continue

        # Safety net: ensure frontmatter is present
        if not has_valid_frontmatter(result):
            logger.warning("  Missing frontmatter — injecting default for: {}", rel_path)
            sources = parse_existing_sources(existing_content) if existing_content else []
            fm = build_frontmatter(page_def, file_path.name, sources)
            result = fm + "\n" + result

        wiki_page_path.write_text(result, encoding="utf-8")
        stats["pages_updated"].append(rel_path)
        logger.info("  → Wrote wiki/{} ({} chars)", rel_path, len(result))

    mark_done(file_path, state)
    save_state(state)

    pages_str = ", ".join(stats["pages_updated"])
    if pages_str:
        append_log(f"ingest | {file_path.name} | {pages_str}")

    stats["status"] = "done" if stats["pages_updated"] else "no_output"
    return stats


# ─────────────────────────────────────────────────────────────────────────────
# Bulk ingest
# ─────────────────────────────────────────────────────────────────────────────

def ingest_all(force: bool = False) -> None:
    if not RAW_DIR.exists():
        console.print(f"[red]raw/ not found at {RAW_DIR}[/red]")
        return

    files = sorted([
        p for p in RAW_DIR.rglob("*")
        if p.is_file() and not should_skip(p)
    ])

    if not files:
        console.print("[yellow]No files found in raw/[/yellow]")
        return

    console.print(
        f"\n[bold]Found {len(files)} files[/bold]  "
        f"[dim]model: {LLM_MODEL}[/dim]\n"
    )

    all_stats = []
    for f in track(files, description="Ingesting..."):
        all_stats.append(ingest_file(f, force=force))

    _print_summary(all_stats)
    rebuild_index()


def _print_summary(all_stats: list[dict]) -> None:
    table = Table(title="Ingest Summary", show_header=True, header_style="bold")
    table.add_column("File",   style="cyan", no_wrap=True, max_width=40)
    table.add_column("Status", style="bold", width=12)
    table.add_column("Pages updated")

    counts: dict[str, int] = {}
    for s in all_stats:
        st = s["status"]
        counts[st] = counts.get(st, 0) + 1
        colour = {
            "done": "green", "skipped": "dim", "empty": "yellow",
            "no_mapping": "red", "no_output": "red",
        }.get(st, "white")
        table.add_row(
            s["file"][:40],
            f"[{colour}]{st}[/{colour}]",
            ", ".join(s["pages_updated"]) or "—",
        )

    console.print(table)
    console.print(
        f"\n[green]Done: {counts.get('done', 0)}[/green]  "
        f"[dim]Skipped: {counts.get('skipped', 0)}[/dim]  "
        f"[yellow]Empty: {counts.get('empty', 0)}[/yellow]  "
        f"[red]No mapping: {counts.get('no_mapping', 0)}[/red]"
    )


# ─────────────────────────────────────────────────────────────────────────────
# Index rebuild — deterministic, no LLM needed
# ─────────────────────────────────────────────────────────────────────────────

def rebuild_index() -> None:
    console.print("\n[bold]Rebuilding index.md...[/bold]")
    pages = sorted([
        p for p in WIKI_DIR.rglob("*.md")
        if p.name not in ("index.md", "log.md", "lint-report.md")
    ])
    if not pages:
        console.print("[yellow]No wiki pages to index.[/yellow]")
        return

    today = date.today().isoformat()
    sections: dict[str, list[tuple[str, str, str]]] = {
        "Operations": [],
        "Fishing":    [],
        "Hunting":    [],
        "Wireless":   [],
    }

    for page in pages:
        rel = str(page.relative_to(WIKI_DIR)).replace("\\", "/")
        content = page.read_text(errors="replace")

        # Pull title from frontmatter
        title = rel
        for line in content.splitlines()[:12]:
            if line.startswith("title:"):
                title = line.replace("title:", "").strip().strip('"\'')
                break

        # First non-frontmatter, non-heading line as summary
        summary = ""
        in_fm = content.strip().startswith("---")
        fm_count = 0
        for line in content.splitlines():
            if line.strip() == "---":
                fm_count += 1
                continue
            if in_fm and fm_count < 2:
                continue
            if line.strip() and not line.startswith("#"):
                summary = line.strip()[:90]
                break

        if rel.startswith("fishing/"):
            sections["Fishing"].append((rel, title, summary))
        elif rel.startswith("hunting/"):
            sections["Hunting"].append((rel, title, summary))
        elif rel.startswith("wireless/"):
            sections["Wireless"].append((rel, title, summary))
        else:
            sections["Operations"].append((rel, title, summary))

    lines = [
        "---",
        f'title: "GSM Outdoors Support Wiki — Index"',
        f"last_updated: {today}",
        f"total_pages: {len(pages)}",
        "---",
        "",
        "# GSM Outdoors Support Wiki — Index",
        "",
    ]

    for section, entries in sections.items():
        if not entries:
            continue
        lines.append(f"## {section}")
        lines.append("")
        lines.append("| Page | What customers find here |")
        lines.append("|---|---|")
        for rel, title, summary in sorted(entries, key=lambda x: x[0]):
            stem = rel.replace(".md", "")
            lines.append(f"| [[{stem}]] — {title} | {summary} |")
        lines.append("")

    INDEX_PATH.write_text("\n".join(lines), encoding="utf-8")
    console.print(f"[green]✓ index.md rebuilt — {len(pages)} pages[/green]")
    append_log(f"reindex | {len(pages)} pages")


# ─────────────────────────────────────────────────────────────────────────────
# Lint
# ─────────────────────────────────────────────────────────────────────────────

def lint_wiki() -> None:
    console.print("\n[bold]Linting wiki...[/bold]")
    # index.md is a navigation file — it intentionally lacks content frontmatter fields
    LINT_SKIP = {"log.md", "lint-report.md", "index.md"}
    pages = [p for p in WIKI_DIR.rglob("*.md")
             if p.name not in LINT_SKIP]
    issues: list[str] = []

    agent_phrases = [
        "tell the customer", "advise the customer", "inform the customer",
        "direct the customer", "the agent should", "instruct the customer",
        "let the customer know", "have the customer",
    ]

    for page in pages:
        rel = str(page.relative_to(WIKI_DIR)).replace("\\", "/")
        content = page.read_text(errors="replace")

        if not has_valid_frontmatter(content):
            issues.append(f"❌ Missing frontmatter: {rel}")

        # Check required frontmatter fields — also accept 'brand:' as alias for 'brands:'
        if has_valid_frontmatter(content):
            end = content.find("\n---\n", 3)
            fm_text = content[3:end] if end != -1 else content[3:]
            for field in ("title:", "topics:", "category:",
                          "source_docs:", "last_updated:", "related_pages:"):
                if field not in fm_text:
                    issues.append(f"⚠️  Missing frontmatter field '{field}': {rel}")
            # brands: or brand: both acceptable
            if "brands:" not in fm_text and "brand:" not in fm_text:
                issues.append(f"⚠️  Missing frontmatter field 'brands:': {rel}")

        lower = content.lower()
        for phrase in agent_phrases:
            if phrase in lower:
                issues.append(f"⚠️  Agent language ('{phrase}'): {rel}")
                break

        lines = content.splitlines()
        half = len(lines) // 2
        if half > 20:
            if "\n".join(lines[:half]).strip()[:200] == \
               "\n".join(lines[half:]).strip()[:200]:
                issues.append(f"⚠️  Duplicate content suspected: {rel}")

        if "⚠️ CONFLICT" in content:
            issues.append(f"⚠️  Unresolved conflict marker: {rel}")

        # Brand isolation check — scan body only, not frontmatter or related_pages links
        page_def = WIKI_PAGES.get(rel)
        if page_def and len(page_def.get("brands", [])) == 1:
            allowed_brand = page_def["brands"][0].lower()

            # Extract body text only (everything after closing ---)
            body_start = content.find("\n---\n", 3)
            body = content[body_start + 5:] if body_start != -1 else content
            # Remove "related_pages" lines and "## Sources" section to avoid false positives
            body_clean = re.sub(r"related_pages:.*", "", body, flags=re.IGNORECASE)
            body_clean = re.sub(r"## Sources.*", "", body_clean, flags=re.DOTALL | re.IGNORECASE)
            body_lower = body_clean.lower()

            # Build rival set — exclude all aliases of the allowed brand
            all_rivals = {
                "phenix", "dobyns", "bonehead", "bucca",
                "avian-x", "avian x", "sog", "walker",
                "connect cellular", "muddy mtrx", "stealth cam",
                "wgi", "wildgame", "boss buck",
            }
            # Remove the allowed brand and any token that appears in it
            allowed_tokens = set(re.split(r"[\s\-]+", allowed_brand))
            rivals = {r for r in all_rivals
                      if r != allowed_brand
                      and not any(tok in r for tok in allowed_tokens if len(tok) > 3)}

            # Per-page exemptions for legitimate cross-brand product name references.
            # Add entries here when a rival brand name appears as part of a product name
            # that is officially sold under the page's own brand.
            # Format: { "relative/path.md": {"rival_word", ...} }
            CONTAMINATION_EXEMPTIONS: dict[str, set[str]] = {
                # The Dobyns MB series is officially named "Mike Bucca BullShad" —
                # it is a Dobyns rod, not a Bucca Brand lure. "bucca" appears here
                # solely as part of the signature rod's proper name.
                "fishing/dobyns-rods.md": {"bucca"},
            }
            page_exemptions = CONTAMINATION_EXEMPTIONS.get(rel, set())

            # for rival in rivals:
            #     if rival in body_lower:
            #         issues.append(
            #             f"⚠️  Possible brand contamination ('{rival}' in body of {rel}"
            #             f" — page should only cover {page_def['brands'][0]})"
            #         )
            #         break
            for rival in rivals:
                if rival in page_exemptions:
                    continue  # legitimate cross-brand product name — skip
                if rival in body_lower:
                    issues.append(
                        f"⚠️  Possible brand contamination ('{rival}' in body of {rel}"
                        f" — page should only cover {page_def['brands'][0]})"
                    )
                    break

    today = date.today().isoformat()
    report_path = WIKI_DIR / "lint-report.md"

    if issues:
        report = f"# Wiki Lint — {today}\n\n{len(issues)} issue(s):\n\n"
        report += "\n".join(issues)
        console.print(f"\n[yellow]{len(issues)} issues:[/yellow]")
        for i in issues:
            console.print(f"  {i}")
    else:
        report = f"# Wiki Lint — {today}\n\n✅ Clean."
        console.print("[green]✓ Wiki is clean[/green]")

    report_path.write_text(report, encoding="utf-8")
    append_log(f"lint | {len(issues)} issues")


# ─────────────────────────────────────────────────────────────────────────────
# Helpers for server.py / auto-ingest API
# ─────────────────────────────────────────────────────────────────────────────

def parse_qa_excel(path: Path) -> dict[str, list[dict]]:
    """
    Parse a Q&A Excel file.
    Returns dict[area_name -> list[{q, a, subject}]].
    Expects columns: Area (or Category), Question, Answer, Subject (optional).
    """
    import pandas as pd

    xl = pd.ExcelFile(str(path))
    by_area: dict[str, list[dict]] = {}

    for sheet in xl.sheet_names:
        try:
            df = xl.parse(sheet, dtype=str).fillna("")
        except Exception:
            continue

        # Normalise column names
        df.columns = [str(c).strip().lower() for c in df.columns]
        col_map: dict[str, str] = {}
        for col in df.columns:
            if col in ("area", "category", "department", "topic"):
                col_map["area"] = col
            elif col in ("question", "q", "issue", "customer question"):
                col_map["q"] = col
            elif col in ("answer", "a", "response", "resolution"):
                col_map["a"] = col
            elif col in ("subject", "title", "summary"):
                col_map["subject"] = col

        if "q" not in col_map or "a" not in col_map:
            continue

        for _, row in df.iterrows():
            area = str(row.get(col_map.get("area", ""), sheet)).strip() or sheet
            q = str(row.get(col_map["q"], "")).strip()
            a = str(row.get(col_map["a"], "")).strip()
            subj = str(row.get(col_map.get("subject", ""), "")).strip()
            if q and a:
                by_area.setdefault(area, []).append({"q": q, "a": a, "subject": subj})

    return by_area


def classify_area_to_wiki(area: str, subject: str, blob: str) -> str:
    """Heuristically map an Excel area name to a wiki page path."""
    text = (area + " " + subject + " " + blob[:200]).lower()
    # Use WIKI_PAGES topics for scoring
    scores: dict[str, int] = {}
    for rel, defn in WIKI_PAGES.items():
        score = 0
        for kw in defn.get("topics", []):
            if kw.lower() in text:
                score += 2
        for brand in defn.get("brands", []):
            if brand.lower() in text:
                score += 3
        scores[rel] = score
    best = max(scores, key=lambda k: scores[k])
    return best if scores[best] > 0 else "fishing/fishing-overview.md"


def default_stub_page(rel_path: str) -> str:
    """Return a minimal stub wiki page for rel_path."""
    page_def = WIKI_PAGES.get(rel_path, {
        "title": rel_path.replace(".md", "").replace("/", " — ").title(),
        "category": rel_path.split("/")[0] if "/" in rel_path else "operations",
        "brands": [], "topics": [], "scope": "", "related": [],
    })
    fm = build_frontmatter(page_def, "")
    return fm + f"\n# {page_def['title']}\n\n*This page is awaiting content from ingest.*\n"


def describe_image(path: Path) -> str:
    """Use Ollama VLM to describe an image."""
    import base64
    try:
        img_b64 = base64.b64encode(path.read_bytes()).decode()
        with httpx.Client(timeout=120) as client:
            resp = client.post(
                f"{OLLAMA_URL}/api/generate",
                json={
                    "model":  LLM_MODEL,
                    "prompt": "Describe this product support image in 2-3 sentences.",
                    "images": [img_b64],
                    "stream": False,
                },
            )
            resp.raise_for_status()
            return resp.json().get("response", "").strip()
    except Exception as e:
        logger.warning("Image description failed: {}", e)
    return ""


def generate_wiki_draft(rel_path: str, qa_pairs: list[dict],
                        existing_content: str) -> str:
    """
    Generate a new wiki section from Q&A pairs (used by server.py auto-ingest).
    Returns the new section text only (not the full page).
    """
    page_def = WIKI_PAGES.get(rel_path, {})
    scope = page_def.get("scope", rel_path)
    brand_rule = _brand_rule(page_def)
    pairs_text = "\n".join(
        f"Q: {qa['q']}\nA: {qa['a']}" for qa in qa_pairs[:20]
    )
    prompt = f"""\
You are writing a new section for a customer-facing support wiki page.
{REFRAME_INSTRUCTION}

TARGET PAGE: wiki/{rel_path}
ABOUT: {scope}
{brand_rule}

Transform these Q&A pairs into clean wiki sections.
Use ## headings, numbered steps for procedures, tables for parts.
Write directly TO the customer — no agent language.
Do NOT include frontmatter — only the body content sections.

Q&A PAIRS TO TRANSFORM:
{pairs_text}

Return ONLY the new wiki sections (no frontmatter, no preamble)."""

    return llm(prompt) or ""


# ─────────────────────────────────────────────────────────────────────────────
# CLI
# ─────────────────────────────────────────────────────────────────────────────

def main() -> None:
    parser = argparse.ArgumentParser(
        description=(
            "GSM Outdoors — Customer-Facing Wiki Builder\n"
            "Transforms agent instruction docs → customer-facing wiki pages.\n"
            "Pages are organised brand-wise within category folders."
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument("--all",     action="store_true", help="Ingest all files in raw/")
    parser.add_argument("--file",    type=Path,           help="Ingest a single file")
    parser.add_argument("--folder",  type=Path,           help="Ingest all files in a folder")
    parser.add_argument("--lint",    action="store_true", help="Run wiki quality checks")
    parser.add_argument("--reindex", action="store_true", help="Rebuild index.md only")
    parser.add_argument("--force",   action="store_true", help="Re-ingest unchanged files")
    args = parser.parse_args()

    # Ensure wiki directory structure exists
    WIKI_DIR.mkdir(parents=True, exist_ok=True)
    for d in ("fishing", "hunting", "wireless"):
        (WIKI_DIR / d).mkdir(exist_ok=True)
    if not LOG_PATH.exists():
        LOG_PATH.write_text("# GSM Outdoors CS Wiki — Log\n\n")

    console.print(
        f"\n[bold cyan]GSM Outdoors — Customer Wiki Builder[/bold cyan]  "
        f"[dim]model: {LLM_MODEL}[/dim]\n"
    )

    if args.lint:
        lint_wiki()

    elif args.reindex:
        rebuild_index()

    elif args.file:
        if not args.file.exists():
            console.print(f"[red]Not found: {args.file}[/red]")
            sys.exit(1)
        stats = ingest_file(args.file, force=args.force)
        console.print(stats)
        rebuild_index()

    elif args.folder:
        if not args.folder.exists():
            console.print(f"[red]Not found: {args.folder}[/red]")
            sys.exit(1)
        files = sorted([
            p for p in args.folder.rglob("*")
            if p.is_file() and not should_skip(p)
        ])
        all_stats = [ingest_file(f, force=args.force)
                     for f in track(files, description="Ingesting...")]
        _print_summary(all_stats)
        rebuild_index()

    elif args.all:
        ingest_all(force=args.force)

    else:
        parser.print_help()


if __name__ == "__main__":
    main()
